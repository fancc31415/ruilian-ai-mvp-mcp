"""
BP 解析模块。
1) extract_text_from_file: 把上传的 PDF / PPTX 转成纯文本
2) parse_bp: 调 LLM 把纯文本结构化成匹配引擎需要的字段
   （没有配置 API Key 时，用关键词规则做一版降级抽取，保证 demo 能跑通）
"""
import re
import pdfplumber
from pptx import Presentation

SECTOR_KEYWORDS = {
    "硬科技": ["硬科技", "半导体材料", "先进制造装备"],
    "半导体": ["半导体", "芯片", "集成电路", "晶圆", "封装测试"],
    "人工智能": ["人工智能", "AI", "大模型", "机器学习", "算法", "AIGC"],
    "医疗健康": ["医疗", "医药", "创新药", "生物", "器械", "健康"],
    "创新药": ["创新药", "生物医药", "临床", "新药研发"],
    "新能源": ["新能源", "光伏", "储能", "锂电", "电池", "氢能"],
    "新消费": ["消费", "品牌", "零售", "食品饮料", "潮玩"],
    "企业服务/SaaS": ["SaaS", "企业服务", "B端软件", "云服务", "ERP", "CRM"],
    "TMT": ["TMT", "互联网", "社交", "内容平台", "媒体"],
    "机器人/智能制造": ["机器人", "智能制造", "自动化", "工业软件"],
    "出海": ["出海", "跨境", "海外市场", "国际化"],
}

STAGE_KEYWORDS = {
    "天使轮": ["天使轮", "种子轮", "Pre-seed", "Seed"],
    "Pre-A轮": ["Pre-A", "天使+"],
    "A轮": ["A轮", "Series A"],
    "B轮": ["B轮", "Series B"],
    "C轮+": ["C轮", "D轮", "Pre-IPO", "Series C", "Series D"],
}


def extract_text_from_file(file_path: str) -> str:
    """按扩展名分流：PDF 用 pdfplumber，PPTX 用 python-pptx。"""
    if file_path.lower().endswith(".pdf"):
        return _extract_pdf(file_path)
    elif file_path.lower().endswith(".pptx"):
        return _extract_pptx(file_path)
    else:
        raise ValueError("仅支持 PDF 或 PPTX 格式的 BP 文件")


def _extract_pdf(file_path: str) -> str:
    chunks = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            chunks.append(text)

    combined = "\n".join(chunks)
    if len(combined.strip()) < 20:  # 提取到的文字太少，大概率是纯图片扫描件，走OCR兜底
        ocr_text = _ocr_pdf(file_path)
        if ocr_text.strip():
            return ocr_text
    return combined


def _ocr_pdf(file_path: str) -> str:
    """对扫描版PDF做OCR兜底：把每页渲染成图片，用tesseract识别中英文文字。
    需要系统装了tesseract-ocr + tesseract-ocr-chi-sim（见项目根目录的 packages.txt）。
    识别失败（比如部署环境没装tesseract）时返回空字符串，上层会提示用户文件无法解析，
    不会假装读到内容。"""
    try:
        import pypdfium2 as pdfium
        import pytesseract

        pdf = pdfium.PdfDocument(file_path)
        chunks = []
        for page in pdf:
            bitmap = page.render(scale=2.0)  # 放大渲染，小字也能识别准一些
            pil_image = bitmap.to_pil()
            text = pytesseract.image_to_string(pil_image, lang="chi_sim+eng")
            chunks.append(text)
        return "\n".join(chunks)
    except Exception:
        return ""


def _extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        chunks.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(chunks)


PARSE_SYSTEM_PROMPT = """你是一名资深FA（财务顾问）行业分析师，擅长快速阅读企业BP（商业计划书）并提炼投融资关键信息。
请仔细阅读用户提供的BP文本，提取以下结构化信息，只返回JSON，不要有多余文字：

{
  "company_name": "公司名称，未提及则填'未提供'",
  "sectors": ["赛道标签，从这些里选1-3个最贴切的：硬科技、半导体、人工智能、医疗健康、创新药、新能源、新消费、企业服务/SaaS、TMT、机器人/智能制造、出海"],
  "stage": "融资阶段，从这些里选一个：天使轮、Pre-A轮、A轮、B轮、C轮+",
  "funding_ask_wan": 本轮融资金额（单位：万元人民币，数字类型，无法判断则填0）,
  "valuation_wan": 估值（单位：万元人民币，数字类型，无法判断则填0）,
  "business_summary": "一句话商业模式概括，30字以内",
  "highlights": ["核心亮点1", "核心亮点2", "核心亮点3"],
  "team_background": "团队背景简述，无信息则填'未提供'",
  "risks_or_gaps": ["BP中信息不足或可能引起机构顾虑的点，最多2条"]
}
"""


def parse_bp(text: str, llm_client) -> dict:
    """优先用 LLM 结构化抽取；LLM 不可用（未配置 Key）时走关键词规则降级方案。"""
    if llm_client.available:
        try:
            result = llm_client.chat_json(PARSE_SYSTEM_PROMPT, text[:12000])
            result.setdefault("sectors", [])
            result.setdefault("stage", "A轮")
            result["_source"] = "llm"
            return result
        except Exception as e:
            fallback = _rule_based_parse(text)
            fallback["_source"] = f"rule_fallback (LLM调用失败: {e})"
            return fallback
    else:
        fallback = _rule_based_parse(text)
        fallback["_source"] = "rule_fallback (未配置API Key)"
        return fallback


def _rule_based_parse(text: str) -> dict:
    """关键词命中规则，作为没有LLM Key时的降级方案，保证功能可跑通。"""
    sectors = []
    for sector, kws in SECTOR_KEYWORDS.items():
        if any(kw in text for kw in kws):
            sectors.append(sector)
    if not sectors:
        sectors = ["TMT"]

    stage = "A轮"
    for s, kws in STAGE_KEYWORDS.items():
        if any(kw in text for kw in kws):
            stage = s
            break

    amount_match = re.search(r"融资.{0,6}?(\d+(?:\.\d+)?)\s*万", text)
    funding_ask = float(amount_match.group(1)) if amount_match else 0

    company_match = re.search(r"([\u4e00-\u9fa5A-Za-z0-9]{2,10})\s*(?:公司|科技|集团)", text)
    company_name = company_match.group(0) if company_match else "未提供"

    return {
        "company_name": company_name,
        "sectors": sectors[:3],
        "stage": stage,
        "funding_ask_wan": funding_ask,
        "valuation_wan": 0,
        "business_summary": text.strip().split("\n")[0][:30] if text.strip() else "未提供",
        "highlights": [],
        "team_background": "未提供",
        "risks_or_gaps": ["未配置LLM，本次为关键词规则粗提取，建议补充手动核对信息"],
    }
